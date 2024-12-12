from pptx import Presentation

# Crear una presentación de PowerPoint
presentation = Presentation()

# Diapositiva 1: Título
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Prácticas Recomendadas por IEEE para la Protección y Coordinación de Sistemas de Energía Industriales y Comerciales"
subtitle.text = ""

# Diapositiva 2: Introducción
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.shapes.placeholders[1]
title.text = "Introducción"
content.text = "¿Qué son las Prácticas Recomendadas por IEEE?\n" \
               "IEEE (Instituto de Ingenieros Eléctricos y Electrónicos) es una organización global que establece " \
               "estándares y prácticas recomendadas para la ingeniería eléctrica.\n" \
               "En el contexto de la protección y coordinación de sistemas de energía, IEEE proporciona lineamientos " \
               "para asegurar la operación segura, confiable y eficiente de sistemas eléctricos industriales y comerciales."

# Diapositiva 3: Importancia de la Protección y Coordinación
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.shapes.placeholders[1]
title.text = "Importancia de la Protección y Coordinación en Sistemas de Energía"
content.text = "Protección: Garantiza que los equipos y componentes eléctricos estén protegidos contra fallas, " \
               "evitando daños costosos y peligrosos.\n" \
               "Coordinación: Asegura que los dispositivos de protección (como interruptores automáticos) actúen " \
               "de manera sincronizada, limitando las interrupciones y minimizando el impacto de las fallas."

# Diapositiva 4: Principales Normas IEEE Relacionadas
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.shapes.placeholders[1]
title.text = "Principales Normas IEEE Relacionadas"
content.text = "1. IEEE C37.2 – Definición de términos para protección, control y automatización de sistemas eléctricos.\n" \
               "2. IEEE C37.112 – Cálculo de la corriente de cortocircuito y el dimensionamiento de interruptores automáticos.\n" \
               "3. IEEE C37.104 – Recomendaciones para la protección de equipos y sistemas eléctricos de distribución.\n" \
               "4. IEEE 1547 – Estándar para la interconexión de recursos energéticos distribuidos."

# Diapositiva 5: Objetivos de la Protección y Coordinación
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.shapes.placeholders[1]
title.text = "Objetivos de la Protección y Coordinación"
content.text = "Seguridad: Proteger a las personas y equipos frente a fallas.\n" \
               "Confiabilidad: Asegurar el suministro continuo de energía en la planta o instalación.\n" \
               "Eficiencia: Optimizar el uso de los dispositivos de protección para minimizar costos operativos."

# Diapositiva 6: Estrategias de Protección
slide_6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.shapes.placeholders[1]
title.text = "Estrategias de Protección"
content.text = "1. Protección Diferencial: Compara la corriente de entrada y salida en un equipo. Si hay una discrepancia, " \
               "indica una falla interna.\n" \
               "2. Protección de Sobrecorriente: Detecta sobrecargas o cortocircuitos en un sistema.\n" \
               "3. Protección de Tierra: Utilizada para detectar fallas a tierra que pueden ser peligrosas."

# Diapositiva 7: Coordinación de Dispositivos de Protección
slide_7 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.shapes.placeholders[1]
title.text = "Coordinación de Dispositivos de Protección"
content.text = "Objetivo: Evitar la desconexión innecesaria de partes del sistema eléctrico.\n" \
               "Métodos Comunes:\n" \
               "1. Coordinación de Tiempo (o de Secuencia): Establece la secuencia en que los dispositivos de protección deben operar.\n" \
               "2. Coordinación de Corriente: Asegura que el dispositivo de protección más cercano a la falla actúe primero."

# Diapositiva 8: Consideraciones en la Protección de Sistemas Industriales y Comerciales
slide_8 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.shapes.placeholders[1]
title.text = "Consideraciones en la Protección de Sistemas Industriales y Comerciales"
content.text = "Industriales: Mayor carga y complejidad debido a equipos de alta potencia.\n"

presentation.save("Presentacion_IEEE.pptx")

